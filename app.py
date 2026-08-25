"""
Redacter — Local PDF/Word/Office redaction tool
Flask backend
"""

from __future__ import annotations
import os
import re
import io
import uuid
import time
import json
import tempfile
import threading
import traceback
from pathlib import Path
from typing import Optional, List

from flask import Flask, request, jsonify, send_file, render_template

# ── optional heavy deps (fail gracefully so server still starts) ──────────────
try:
    import fitz  # pymupdf
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

try:
    from docx import Document as DocxDocument
    from docx.shared import RGBColor
    from docx.enum.text import WD_COLOR_INDEX
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from PIL import Image, ImageFilter, ImageDraw
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

try:
    import pytesseract
    HAS_TESSERACT = True
except ImportError:
    HAS_TESSERACT = False

# ── Flask app ──────────────────────────────────────────────────────────────────
app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = None   # no file size limit

BASE_DIR = Path(__file__).parent
IS_CLOUD = os.environ.get('RAILWAY_ENVIRONMENT') or os.environ.get('RENDER')
if IS_CLOUD:
    UPLOAD_DIR = Path(tempfile.gettempdir()) / 'redacter_uploads'
    OUTPUT_DIR = Path(tempfile.gettempdir()) / 'redacter_outputs'
else:
    UPLOAD_DIR = BASE_DIR / 'uploads'
    OUTPUT_DIR = BASE_DIR / 'outputs'
UPLOAD_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)

# in-memory job store  {file_id: {path, original_name, ext, uploaded_at}}
FILE_STORE: dict = {}
# cleanup lock
_store_lock = threading.Lock()

# ── Auto-detect regex patterns ─────────────────────────────────────────────────
PATTERNS = {
    "dates": [
        r'\b\d{1,2}/\d{1,2}/\d{2,4}\b',
        r'\b\d{4}-\d{2}-\d{2}\b',
        r'\b(?:January|February|March|April|May|June|July|August|September|October|November|December|Jan|Feb|Mar|Apr|Jun|Jul|Aug|Sep|Oct|Nov|Dec)\.?\s+\d{1,2},?\s+\d{4}\b',
        r'\b\d{1,2}\s+(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}\b',
    ],
    "times": [
        r'\b\d{1,2}:\d{2}(?::\d{2})?\s*(?:AM|PM|am|pm|a\.m\.|p\.m\.)?\b',
    ],
    "phones": [
        r'\b(?:\+?1[\s.\-]?)?\(?\d{3}\)?[\s.\-]?\d{3}[\s.\-]?\d{4}\b',
    ],
    "emails": [
        r'\b[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}\b',
    ],
    "ssn": [
        r'\b\d{3}[-\s]?\d{2}[-\s]?\d{4}\b',
        r'\b[A-Z]{1,2}\d{6,9}\b',
    ],
}

NAME_STOPWORDS = {
    # Common doc words
    'The','This','That','With','From','Page','Date','Time','Section',
    'Table','Figure','Dear','Kind','Best','Regards','Hello','Subject',
    'Re','Cc','Bcc','To','In','At','On','Of','For','And','Or','But',
    # Inspection/report words
    'Visual','Testing','Remote','Drain','Weld','Label','Observation',
    'Comments','Applied','Technical','Services','Inspection','General',
    'Main','Scope','Access','Findings','Signoff','Report','Summary',
    'Section','Appendix','Total','Pass','Fail','Result','Results',
    'Inspector','Level','Certificate','Standard','Method','Procedure',
    'Client','Customer','Project','Location','Facility','Plant','Site',
    'Drawing','Reference','Number','Rev','Page','Of','Per','Per',
    'Note','Notes','See','Refer','Attached','Enclosure','Exhibit',
    'January','February','March','April','June','July','August',
    'September','October','November','December',
    # Single-letter second names (like 'Drain A', 'Weld B')
}


def find_names(text: str) -> List[str]:
    # Require BOTH words >= 3 chars to reduce false positives
    matches = re.findall(r'\b([A-Z][a-z]{2,20})\s+([A-Z][a-z]{2,20})\b', text)
    seen = set()
    result = []
    for f, l in matches:
        name = f"{f} {l}"
        if (f not in NAME_STOPWORDS
                and l not in NAME_STOPWORDS
                and name not in seen):
            seen.add(name)
            result.append(name)
    return result


def gather_terms(options: dict) -> List[str]:
    """Build final list of strings to redact from UI options dict."""
    manual       = options.get('manual_terms', [])
    case_sensitive = options.get('case_sensitive', False)

    # collect text from the file for auto-detect (caller must pass 'full_text')
    full_text = options.get('_full_text', '')
    terms = list(manual)

    auto = options.get('auto_detect', {})

    def add_pattern_matches(keys):
        for key in keys:
            for pat in PATTERNS.get(key, []):
                for m in re.finditer(pat, full_text, re.IGNORECASE):
                    t = m.group(0)
                    if t not in terms:
                        terms.append(t)

    if auto.get('names'):
        for n in find_names(full_text):
            if n not in terms:
                terms.append(n)
    if auto.get('dates'):  add_pattern_matches(['dates'])
    if auto.get('times'):  add_pattern_matches(['times'])
    if auto.get('phones'): add_pattern_matches(['phones'])
    if auto.get('emails'): add_pattern_matches(['emails'])
    if auto.get('ssn'):    add_pattern_matches(['ssn'])

    return terms


# ── Text extraction helpers ────────────────────────────────────────────────────

def extract_pdf_text(path: str, use_ocr: bool = False) -> tuple:
    """Return (full_text, image_regions).
       image_regions: list of {page, xref, rect, pil_image, text}
    """
    if not HAS_FITZ:
        return '', []
    doc = fitz.open(path)
    pages_text = []
    image_regions = []
    for page_num, page in enumerate(doc):
        pages_text.append(page.get_text())
        if use_ocr and HAS_TESSERACT and HAS_PIL:
            for img in page.get_images(full=True):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    pil_img = Image.open(io.BytesIO(base_image['image'])).convert('RGB')
                    ocr_text = pytesseract.image_to_string(pil_img)
                    rects = page.get_image_rects(xref)
                    if rects and ocr_text.strip():
                        image_regions.append({
                            'page': page_num,
                            'xref': xref,
                            'rect': rects[0],
                            'pil_image': pil_img,
                            'text': ocr_text,
                        })
                except Exception:
                    pass
    doc.close()
    return '\n'.join(pages_text), image_regions


def extract_docx_text(path: str) -> str:
    if not HAS_DOCX: return ''
    doc = DocxDocument(path)
    parts = [p.text for p in doc.paragraphs]
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return '\n'.join(parts)


def extract_xlsx_text(path: str) -> str:
    if not HAS_OPENPYXL: return ''
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    parts.append(str(cell))
    wb.close()
    return '\n'.join(parts)


def extract_pptx_text(path: str) -> str:
    if not HAS_PPTX: return ''
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        parts.append(run.text)
    return '\n'.join(parts)


def extract_text(path: str, ext: str, use_ocr: bool = False) -> tuple:
    """Return (text, image_regions). image_regions only relevant for PDFs."""
    if ext == '.pdf':
        return extract_pdf_text(path, use_ocr)
    elif ext in ('.docx', '.doc'):
        return extract_docx_text(path), []
    elif ext == '.xlsx':
        return extract_xlsx_text(path), []
    elif ext == '.pptx':
        return extract_pptx_text(path), []
    elif ext == '.txt':
        return Path(path).read_text(errors='replace'), []
    elif ext in ('.png', '.jpg', '.jpeg', '.tiff', '.tif'):
        if HAS_TESSERACT and HAS_PIL:
            img = Image.open(path).convert('RGB')
            return pytesseract.image_to_string(img), []
        return '', []
    return '', []


# ── Redaction helpers ──────────────────────────────────────────────────────────

def build_regex(terms: List[str], case_sensitive: bool = False) -> Optional[re.Pattern]:
    if not terms:
        return None
    flags = 0 if case_sensitive else re.IGNORECASE
    pattern = '|'.join(re.escape(t) for t in sorted(terms, key=len, reverse=True))
    return re.compile(pattern, flags)


def redact_pdf(src: str, dst: str, terms: List[str], options: dict):
    if not HAS_FITZ:
        raise RuntimeError("pymupdf not installed")
    style       = options.get('style', 'black')
    case_sens   = options.get('case_sensitive', False)
    use_ocr     = options.get('use_ocr', False)
    regex       = build_regex(terms, case_sens)

    doc = fitz.open(src)
    for page in doc:
        if regex:
            # text layer redaction
            for m in regex.finditer(page.get_text()):
                # find rects for every occurrence on the page
                pass
            # use fitz search for accurate rects
            for term in terms:
                flags = 0 if case_sens else fitz.TEXT_PRESERVE_WHITESPACE
                rects = page.search_for(term, quads=False)
                for rect in rects:
                    if style == 'blur':
                        # render region, blur with Pillow, re-insert
                        _blur_pdf_rect(page, rect)
                    else:
                        annot = page.add_redact_annot(rect, fill=(0, 0, 0))
            page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)

        # OCR-based image redaction — parallel per page
        if use_ocr and HAS_TESSERACT and HAS_PIL and regex:
            page_images = []
            for img in page.get_images(full=True):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    raw = base_image['image']
                    if len(raw) < 20000:
                        continue
                    pil_img = Image.open(io.BytesIO(raw)).convert('RGB')
                    if pil_img.width < 100 or pil_img.height < 100:
                        continue
                    page_images.append((xref, pil_img))
                except Exception:
                    pass

            def _process_one_image(args):
                xref, pil_img = args
                try:
                    # Resize for speed
                    if pil_img.width > 1200:
                        ratio = 1200 / pil_img.width
                        pil_img = pil_img.resize(
                            (1200, int(pil_img.height * ratio)), Image.LANCZOS)
                    # Corner-region OCR: timestamps always in corners
                    # Scan top-left, top-right, bottom-left, bottom-right (15% of image)
                    w, h = pil_img.width, pil_img.height
                    cx, cy = int(w * 0.35), int(h * 0.15)
                    corners = [
                        (0, 0, cx, cy),           # top-left
                        (w-cx, 0, w, cy),          # top-right
                        (0, h-cy, cx, h),          # bottom-left
                        (w-cx, h-cy, w, h),        # bottom-right
                    ]
                    corner_has_match = False
                    for box in corners:
                        crop = pil_img.crop(box)
                        try:
                            txt = pytesseract.image_to_string(crop, config='--psm 6')
                            if regex.search(txt):
                                corner_has_match = True
                                break
                        except Exception:
                            pass
                    # Full redact pass only if corner matched OR do full scan
                    if corner_has_match or True:  # always do full for accuracy
                        pil_img = _redact_image_with_ocr(pil_img, terms, case_sens, style)
                    img_bytes = io.BytesIO()
                    pil_img.save(img_bytes, format='JPEG', quality=90)
                    return xref, img_bytes.getvalue()
                except Exception:
                    return xref, None

            from concurrent.futures import ThreadPoolExecutor
            with ThreadPoolExecutor(max_workers=4) as pool:
                results = list(pool.map(_process_one_image, page_images))
            for xref, data in results:
                if data:
                    try:
                        doc.update_image(xref, stream=data)
                    except Exception:
                        pass

    doc.save(dst, garbage=4, deflate=True)
    doc.close()


def _blur_pdf_rect(page, rect):
    """Render a PDF rect, blur it, then draw it back as an image."""
    try:
        clip = fitz.Rect(rect)
        mat = fitz.Matrix(3, 3)  # 3× zoom for quality
        pix = page.get_pixmap(matrix=mat, clip=clip, alpha=False)
        pil = Image.open(io.BytesIO(pix.tobytes('png'))).convert('RGB')
        pil = pil.filter(ImageFilter.GaussianBlur(radius=12))
        buf = io.BytesIO()
        pil.save(buf, format='PNG')
        buf.seek(0)
        img_rect = fitz.Rect(rect)
        page.insert_image(img_rect, stream=buf.read(), overlay=True)
    except Exception:
        # fallback: grey box
        page.draw_rect(rect, color=(0.5, 0.5, 0.5), fill=(0.5, 0.5, 0.5))


def _redact_image_with_ocr(pil_img: 'Image.Image', terms: List[str],
                            case_sensitive: bool, style: str) -> 'Image.Image':
    """Find text regions in a PIL image via Tesseract and black-out/blur them."""
    try:
        data = pytesseract.image_to_data(pil_img, output_type=pytesseract.Output.DICT)
        draw = ImageDraw.Draw(pil_img)
        flags = 0 if case_sensitive else re.IGNORECASE
        pattern = build_regex(terms, case_sensitive)
        if not pattern:
            return pil_img

        n = len(data['text'])
        for i in range(n):
            word = data['text'][i]
            if not word.strip():
                continue
            if pattern.search(word):
                x, y, w, h = data['left'][i], data['top'][i], data['width'][i], data['height'][i]
                if style == 'blur':
                    region = pil_img.crop((x, y, x+w, y+h))
                    region = region.filter(ImageFilter.GaussianBlur(radius=10))
                    pil_img.paste(region, (x, y))
                else:
                    draw.rectangle([x, y, x+w, y+h], fill=(0, 0, 0))
    except Exception:
        pass
    return pil_img


def redact_docx(src: str, dst: str, terms: List[str], options: dict):
    if not HAS_DOCX:
        raise RuntimeError("python-docx not installed")
    regex = build_regex(terms, options.get('case_sensitive', False))
    if not regex:
        import shutil; shutil.copy2(src, dst); return

    doc = DocxDocument(src)

    def redact_paragraph(para):
        # rebuild runs to avoid splitting mid-word
        full = para.text
        if not regex.search(full):
            return
        redacted = regex.sub(lambda m: '█' * len(m.group(0)), full)
        # clear runs and set plain text
        for run in para.runs:
            run.text = ''
        if para.runs:
            para.runs[0].text = redacted
            para.runs[0].font.color.rgb = RGBColor(0, 0, 0)
        else:
            run = para.add_run(redacted)
            run.font.color.rgb = RGBColor(0, 0, 0)

    for para in doc.paragraphs:
        redact_paragraph(para)
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    redact_paragraph(para)

    doc.save(dst)


def redact_xlsx(src: str, dst: str, terms: List[str], options: dict):
    if not HAS_OPENPYXL:
        raise RuntimeError("openpyxl not installed")
    regex = build_regex(terms, options.get('case_sensitive', False))
    if not regex:
        import shutil; shutil.copy2(src, dst); return

    wb = openpyxl.load_workbook(src)
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    if regex.search(cell.value):
                        cell.value = regex.sub(lambda m: '██████', cell.value)
    wb.save(dst)
    wb.close()


def redact_pptx(src: str, dst: str, terms: List[str], options: dict):
    if not HAS_PPTX:
        raise RuntimeError("python-pptx not installed")
    regex = build_regex(terms, options.get('case_sensitive', False))
    if not regex:
        import shutil; shutil.copy2(src, dst); return

    prs = Presentation(src)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if regex.search(run.text):
                            run.text = regex.sub(lambda m: '█' * len(m.group(0)), run.text)
    prs.save(dst)


def redact_txt(src: str, dst: str, terms: List[str], options: dict):
    regex = build_regex(terms, options.get('case_sensitive', False))
    text  = Path(src).read_text(errors='replace')
    if regex:
        text = regex.sub(lambda m: '█' * len(m.group(0)), text)
    Path(dst).write_text(text)


def redact_image(src: str, dst: str, terms: List[str], options: dict):
    if not (HAS_PIL and HAS_TESSERACT):
        raise RuntimeError("Pillow/pytesseract not installed")
    style = options.get('style', 'black')
    img = Image.open(src).convert('RGB')
    img = _redact_image_with_ocr(img, terms, options.get('case_sensitive', False), style)
    img.save(dst)


def do_redact(file_id: str, options: dict) -> str:
    """Core redaction dispatcher. Returns output path."""
    with _store_lock:
        info = FILE_STORE.get(file_id)
    if not info:
        raise ValueError(f"Unknown file_id: {file_id}")

    src   = info['path']
    ext   = info['ext']
    name  = info['original_name']
    stem  = Path(name).stem.replace(' ','_').replace('%','_')  # sanitize spaces
    out_name = f"{stem}_REDACTED{ext}"
    dst   = str(OUTPUT_DIR / f"{file_id}_{out_name}")

    use_ocr = options.get('use_ocr', False)
    full_text, _ = extract_text(src, ext, use_ocr)
    options['_full_text'] = full_text

    # Use direct terms from JS matches array if provided, otherwise detect
    terms = options.get('_direct_terms') or gather_terms(options)

    if ext == '.pdf':
        redact_pdf(src, dst, terms, options)
    elif ext in ('.docx', '.doc'):
        redact_docx(src, dst, terms, options)
    elif ext == '.xlsx':
        redact_xlsx(src, dst, terms, options)
    elif ext == '.pptx':
        redact_pptx(src, dst, terms, options)
    elif ext == '.txt':
        redact_txt(src, dst, terms, options)
    elif ext in ('.png', '.jpg', '.jpeg', '.tiff', '.tif'):
        redact_image(src, dst, terms, options)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    return dst


# ── Cleanup thread ─────────────────────────────────────────────────────────────
def _cleanup_old_files():
    while True:
        time.sleep(60)
        cutoff = time.time() - 600  # 10 min
        with _store_lock:
            stale = [fid for fid, info in FILE_STORE.items()
                     if info.get('uploaded_at', 0) < cutoff]
            for fid in stale:
                try:
                    os.unlink(FILE_STORE[fid]['path'])
                except Exception:
                    pass
                del FILE_STORE[fid]
        # also clean output dir
        for p in OUTPUT_DIR.glob('*'):
            if p.stat().st_mtime < cutoff:
                try:
                    p.unlink()
                except Exception:
                    pass

threading.Thread(target=_cleanup_old_files, daemon=True).start()


# ── Routes ─────────────────────────────────────────────────────────────────────

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/upload', methods=['POST'])
def upload():
    if 'file' not in request.files:
        return jsonify(error='No file part'), 400
    f = request.files['file']
    if not f.filename:
        return jsonify(error='No filename'), 400

    ext = Path(f.filename).suffix.lower()
    allowed = {'.pdf','.docx','.doc','.xlsx','.pptx','.txt',
               '.png','.jpg','.jpeg','.tiff','.tif'}
    if ext not in allowed:
        return jsonify(error=f'Unsupported file type: {ext}'), 400

    file_id = str(uuid.uuid4())
    save_path = str(UPLOAD_DIR / f"{file_id}{ext}")
    f.save(save_path)

    # Count pages for PDF
    page_count = 1
    if ext == '.pdf':
        try:
            import fitz
            doc = fitz.open(save_path)
            page_count = len(doc)
            doc.close()
        except Exception:
            page_count = 1

    with _store_lock:
        FILE_STORE[file_id] = {
            'path': save_path,
            'original_name': f.filename,
            'ext': ext,
            'page_count': page_count,
            'uploaded_at': time.time(),
        }
    return jsonify(
        file_id=file_id,
        filename=f.filename,
        ext=ext,
        page_count=page_count,
    )


@app.route('/preview', methods=['POST'])
def preview():
    data    = request.get_json(force=True)
    file_id = data.get('file_id')
    options = data.get('options', {})

    with _store_lock:
        info = FILE_STORE.get(file_id)
    if not info:
        return jsonify(error='Unknown file_id'), 404

    try:
        full_text, _ = extract_text(info['path'], info['ext'],
                                    options.get('use_ocr', False))
        options['_full_text'] = full_text
        terms = gather_terms(options)

        # find actual matches with context
        matches = []
        if terms:
            regex = build_regex(terms, options.get('case_sensitive', False))
            for m in regex.finditer(full_text):
                start = max(0, m.start() - 40)
                end   = min(len(full_text), m.end() + 40)
                ctx   = full_text[start:end].replace('\n', ' ')
                matches.append({
                    'term': m.group(0),
                    'context': f"…{ctx}…",
                })
                if len(matches) >= 200:
                    break

        return jsonify(terms=terms, matches=matches, total=len(matches))
    except Exception as e:
        return jsonify(error=str(e), traceback=traceback.format_exc()), 500


# ── Async job store ───────────────────────────────────────────────────────────
JOB_STORE: dict = {}   # job_id -> {status, dst, error, out_name}
_job_lock = threading.Lock()

def _run_redact_job(job_id: str, file_id: str, options: dict):
    """Background thread: runs redaction and updates JOB_STORE."""
    try:
        dst = do_redact(file_id, options)
        out_name = Path(dst).name.split('_', 1)[-1] if '_' in Path(dst).name else Path(dst).name
        with _job_lock:
            JOB_STORE[job_id] = {'status': 'done', 'dst': dst, 'out_name': out_name}
    except Exception as e:
        with _job_lock:
            JOB_STORE[job_id] = {'status': 'error', 'error': str(e),
                                  'traceback': traceback.format_exc()}


@app.route('/redact', methods=['POST'])
def redact():
    data     = request.get_json(force=True)
    file_id  = data.get('file_id')
    options  = data.get('options', {})
    matches_list = data.get('matches', [])
    style    = data.get('style', options.get('style', 'black'))
    options['style'] = style

    if matches_list:
        direct_terms = list(dict.fromkeys(
            m['text'] for m in matches_list if m.get('text')
        ))
        if direct_terms:
            options['_direct_terms'] = direct_terms

    # Kick off background job — return immediately so HTTP doesn't time out
    job_id = str(uuid.uuid4())
    with _job_lock:
        JOB_STORE[job_id] = {'status': 'processing'}

    t = threading.Thread(target=_run_redact_job, args=(job_id, file_id, options), daemon=True)
    t.start()

    return jsonify(job_id=job_id), 202


@app.route('/job/<job_id>')
def job_status(job_id: str):
    with _job_lock:
        job = JOB_STORE.get(job_id)
    if not job:
        return jsonify(error='Unknown job'), 404
    if job['status'] == 'done':
        return jsonify(status='done', download_url=f'/download/{Path(job["dst"]).name}',
                       out_name=job['out_name'])
    if job['status'] == 'error':
        return jsonify(status='error', error=job.get('error', 'Unknown error'))
    return jsonify(status='processing')


@app.route('/download/<path:dl_id>')
def download(dl_id: str):
    from urllib.parse import unquote
    safe = Path(unquote(dl_id)).name  # URL-decode + strip any path
    p = OUTPUT_DIR / safe
    if not p.exists():
        return jsonify(error='File not found'), 404
    return send_file(str(p), as_attachment=True, download_name=safe.split('_', 1)[-1] if '_' in safe else safe)


@app.route('/status')
def status():
    return jsonify(
        ok=True,
        has_fitz=HAS_FITZ,
        has_docx=HAS_DOCX,
        has_openpyxl=HAS_OPENPYXL,
        has_pptx=HAS_PPTX,
        has_pil=HAS_PIL,
        has_tesseract=HAS_TESSERACT,
    )


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5050))
    app.run(host='0.0.0.0', port=port, debug=False, threaded=True)
